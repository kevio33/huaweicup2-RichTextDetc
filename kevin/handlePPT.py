from pptx import Presentation
import aspose.slides as slides
import aspose.pydrawing as drawing
import os

# import sys
# sys.path.append(r'..')
from utils.fileInfo import FileInfo
from kevin.handlePic import OCR
from regSensitive import regexSensitive

'''
    读取ppt内容
    只能读取pptx内容，而无法读取ppt
    因此要先转换格式
'''
def readPPt(fileName,filePath,isPPT=False):
    pptx = None
    if isPPT == True:
        # 转换为pptx格式
        with slides.Presentation(filePath) as presentation:
            filePath = filePath[:-4]
            # print("\n"+filePath)
            filePath = "{0}.pptx".format(filePath)
            presentation.save(filePath, slides.export.SaveFormat.PPTX)

    pptx = Presentation(filePath)# 文件对象
    textRest = []# 用来存放每一页文本内容
    for slide in pptx.slides: # 每个幻灯片
        for shape in slide.shapes: #幻灯片的所有形状
            if shape.has_text_frame: #如果有文本框
                text_frame = shape.text_frame 
                text = text_frame.text # 直接获取文本框的所有内容，这里因为要考虑相关性，所有不更具体导paragraph和run
                textRest.append(text)
            else:
                # 抓取图片，将图片保存到'kevin/{fileName}_img/'目录下面    
                try:
                    imdata = shape.image.blob
                    # 判断文件后缀类型
                    imagetype = shape.image.content_type
                    typekey = imagetype.find('/') + 1
                    imtype = imagetype[typekey:]
                    # 创建image文件夹保存抽出图片
                    rootPath = os.path.abspath(__file__)
                    rootPath = rootPath[:-12] #父目录
                    path = rootPath+"{}_image/".format(fileName)
                    if not os.path.exists(path):
                        os.makedirs(path)
                    # 图片生成
                    image_file = path + shape.name + "." + imtype
                    # name = shape.name
                    file_str=open(image_file,'wb')
                    file_str.write(imdata)
                    file_str.close()
                    # imgList.append(FileInfo(fileName=name,filePath=image_file))
                    res = OCR(image_file) # 读取图片内容
                    textRest.extend(res)
                except:
                    pass
    return textRest

def handlePpt(fileName,filePath,isPPT=False):
    
    textRes = readPPt(fileName,filePath,isPPT)
    textRes = regexSensitive(textRes)
    out_file = open('kevin/{}.txt'.format(fileName),mode='w',encoding='utf-8') #保存路径要注意，如果是main调用进来需要改变路径
    out_file.write(str(textRes))
    out_file.close()
    

if __name__ == '__main__':
    import sys
    sys.path.append(r'..')
    from utils.fileInfo import FileInfo

    lis = handlePpt('20180327081403010127.ppt',r'E:\huaweicup\huaweicup2-RichTextDetc\赛题材料\office\20180327081403010127.ppt',isPPT=True)
    print(str(lis))
    